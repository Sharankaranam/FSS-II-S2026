from pptx import Presentation

def extract_text(filename):
    print(f"--- {filename} ---")
    try:
        prs = Presentation(filename)
        for i, slide in enumerate(prs.slides):
            print(f"Slide {i+1}:")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    print("  Text:", shape.text.strip())
            print()
    except Exception as e:
        print(f"Error reading {filename}: {e}")

extract_text('fuzzy.pptx')
