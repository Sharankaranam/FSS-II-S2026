from pptx import Presentation
import sys

def replace_text(shape, old_text, new_text):
    if not shape.has_text_frame:
        return
    # The text can be split across multiple runs
    # To replace substring safely over split runs, replace it cleanly if found entirely in a run, or we read all text.
    # Python-pptx handles text in chunks "runs"
    # An easier way is if old_text is just matched inside the whole shape text:
    if old_text in shape.text:
        # Simplistic approach: just clear the frame and add the new text if it's a simple shape, but it destroys formatting.
        # So we just scan across runs with simple matches.
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)

try:
    prs = Presentation('fuzzy.pptx')
    
    # 1. Slide 6 (index 5)
    for shape in prs.slides[5].shapes:
        replace_text(shape, 'triangular/trapezoidal', 'triangular')
        
    # 2. Slide 7 (index 6)
    for shape in prs.slides[6].shapes:
        replace_text(shape, 'Fuzzy Rule Base & Rule Evaluation', 'Fuzzy Inference & Dynamic Priority Engine')
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if 'Instead of checking all 243 combinations' in r.text or 'optimized rules' in r.text:
                        r.text = r.text.replace('Instead of checking all 243 combinations, we use optimized rules (~40–50 rules)', 'Instead of rigid multi-condition matrices, the system individually maps parameters to Zero-Order constants.')
                    if 'operator to determine the rule strength' in r.text or 'The AND condition' in r.text:
                        r.text = r.text.replace('The AND condition in each rule is evaluated using the MIN operator to determine the rule strength.', "The engine applies dynamic weights scaled by the user's physical needs directly to the independent parameter scores.")
                    if 'IF Taste' in r.text:
                        r.text = "Each parameter contributes uniquely before being multiplied by context weights."
                    
    # 3. Slide 10 (index 9)
    for shape in prs.slides[9].shapes:
        replace_text(shape, 'GeoJSON', 'CSV Custom Parsers')
        replace_text(shape, 'Geolocation APIs', 'Vanilla JS Engine')

    # 4. Add a new slide about Dynamic Priority Engine
    title_slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "The Dynamic Context Priority Engine"
    tf = None
    for shape in slide.placeholders:
        if shape.has_text_frame and shape.name.find('Title') == -1:
            tf = shape.text_frame
            break
    
    if not tf:
        from pptx.util import Inches
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = txBox.text_frame

    p = tf.paragraphs[0]
    p.text = "Traditional recommendation systems are rigid. To combat this, we built a physical context heuristic interceptor."
    
    p2 = tf.add_paragraph()
    p2.text = "The UI captures 5 real-time states (Hunger, Fatigue, Restroom Urgency, Hygiene, and Budget)."
    
    p3 = tf.add_paragraph()
    p3.text = "The Engine mathematically intercepts these states. For example, if a user flags 'Fatigue', the engine calculates an artificial boost multiplier and heavily increases the importance of the 'Distance' parameter."

    p4 = tf.add_paragraph()
    p4.text = "This allows the fuzzy logic to behave dynamically, solving problems intuitively based on the traveler's exact state of exhaustion."

    # Move slide
    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(7, slides[-1])

    prs.save('fuzzy_updated.pptx')
    print("SUCCESS: Saved as fuzzy_updated.pptx")
except Exception as e:
    import traceback
    traceback.print_exc()
